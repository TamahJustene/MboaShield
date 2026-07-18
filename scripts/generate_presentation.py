#!/usr/bin/env python3
"""Generate the MboaShield SIN 2026 presentation deck."""

from __future__ import annotations

from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations"
STATIC_OUT = ROOT / "frontend" / "static" / "presentations"
FILENAME = "MboaShield_SIN2026.pptx"

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover
    raise SystemExit("Install python-pptx: pip install python-pptx") from exc

SLIDES: list[tuple[str, list[str]]] = [
    (
        "MboaShield",
        [
            "Sovereign AI shield against deepfakes, disinformation and digital identity theft",
            "Made in Cameroon - SIN 2026 National Best ICT Project",
            "Solo founder: Justene Nkwagoh Tamah",
            "Theme: Protect cyberspace from AI excesses and digital patriotism",
        ],
    ),
    (
        "The problem (Cameroon, today)",
        [
            "Fake minister announcements and curfew hoaxes on WhatsApp",
            "Voice clones demanding MoMo transfers",
            "Deepfake images targeting institutions and public figures",
            "Youth forward before verifying - fraud, panic, eroded trust",
        ],
    ),
    (
        "The solution",
        [
            "Detect synthetic media and scam patterns",
            "Verify rumours and flag fake official accounts",
            "Train Mboa Ambassadors in digital patriotism",
            "WhatsApp-style demo narrative - partial FR/EN UI - privacy-aware sovereign stack",
        ],
    ),
    (
        "Live demo (90 seconds)",
        [
            "1. Viral rumour - risk score and source check",
            "2. Fake MINPOSTEL-style account - impersonation",
            "3. Minister voice sample - clone risk",
            "4. Suspicious image - synthetic signals",
            "5. Mboa Ambassador certificate",
            "Demo: mboashield.onrender.com - Run Grand Jury demo",
        ],
    ),
    (
        "National platform depth (v2.8.0)",
        [
            "TrustAssessment: one explainable object across text, identity, audio, image and intelligence",
            "Human-review incident workflow, NTOC, evidence and signed official communications",
            "Institution trust network plus webhooks, STIX, CAP, TAXII 2.1 and SCIM 2.0",
            "Governance, country packs, sector modules, HA/DR and load-test patterns",
        ],
    ),
    (
        "Innovation and AI (honest)",
        [
            "Text, audio, image, identity, civic modules - multimodal trust engine",
            "17-institution registry for local impersonation detection",
            "Heuristic and registry today; ONNX path with checksums for production",
            "Differentiator: Cameroon WhatsApp context, institutions, civic mission",
        ],
    ),
    (
        "Impact and digital patriotism",
        [
            "Citizens: free checks and education",
            "Schools and youth: Mboa Ambassadors",
            "Media, banks, telcos: B2B API",
            "Government: registry, verified comms, national analytics",
        ],
    ),
    (
        "Business model",
        [
            "Citizen Check: free (mission and sponsors)",
            "Pro: media/SMEs - 75k to 250k FCFA/month",
            "Enterprise: banks/telcos - annual API",
            "Ambassadors: schools - training packs",
            "Year 1 target: pilots and 18-25M FCFA revenue path",
        ],
    ),
    (
        "The ask",
        [
            "Special Prize of the President of the Republic to:",
            "1. Industrialise MboaShield nationally",
            "2. Launch institutional identity registry v1 at scale",
            "3. Train 50,000 Mboa Ambassadors in Year 1",
            "MboaShield IS the 2026 theme - not adjacent to it.",
            "tamahjustene45@gmail.com - github.com/TamahJustene/MboaShield",
        ],
    ),
]

INK = RGBColor(7, 31, 26)
EMERALD = RGBColor(34, 197, 94)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(203, 213, 225)


def _style_slide(slide, number: int) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = INK

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = EMERALD
    accent.line.fill.background()

    footer = slide.shapes.add_textbox(
        Inches(0.72), Inches(7.08), Inches(11.9), Inches(0.24)
    )
    paragraph = footer.text_frame.paragraphs[0]
    paragraph.text = f"MBOASHIELD · SOVEREIGN DIGITAL TRUST INFRASTRUCTURE                                      {number:02d}"
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = MUTED
    paragraph.alignment = PP_ALIGN.LEFT


def _add_bullet_slide(
    prs: Presentation, title: str, bullets: list[str], number: int
) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _style_slide(slide, number)
    slide.shapes.title.text = title
    title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
    title_run.font.name = "Aptos Display"
    title_run.font.size = Pt(30)
    title_run.font.bold = True
    title_run.font.color.rgb = WHITE
    body = slide.placeholders[1].text_frame
    body.clear()
    body.margin_left = Inches(0.12)
    first = True
    for line in bullets:
        if not line.strip():
            continue
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(21)
        p.font.color.rgb = WHITE
        p.space_after = Pt(13)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    _style_slide(title_slide, 1)
    title_slide.shapes.title.text = SLIDES[0][0]
    title_slide.placeholders[1].text = "\n".join(SLIDES[0][1])
    title_paragraph = title_slide.shapes.title.text_frame.paragraphs[0]
    title_paragraph.font.name = "Aptos Display"
    title_paragraph.font.size = Pt(48)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = WHITE
    subtitle = title_slide.placeholders[1].text_frame
    for paragraph in subtitle.paragraphs:
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = MUTED
        paragraph.space_after = Pt(6)

    for number, (title, bullets) in enumerate(SLIDES[1:], start=2):
        _add_bullet_slide(prs, title, bullets, number)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    STATIC_OUT.mkdir(parents=True, exist_ok=True)
    path_docs = OUT_DIR / FILENAME
    path_static = STATIC_OUT / FILENAME
    prs.save(path_docs)
    prs.save(path_static)
    return path_docs


def main() -> int:
    path = build()
    print(f"Wrote {path}")
    print(f"Wrote {STATIC_OUT / FILENAME}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
